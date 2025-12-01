# Enhanced PPTX Reader using LlamaIndex approach
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional, Union
from pptx import Presentation

from llama_index.core.readers.base import BaseReader
from llama_index.core.schema import Document

logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s - %(levelname)s - %(message)s",
)
logger = logging.getLogger(__name__)


class PptxReader(BaseReader):
    """
    Enhanced PowerPoint parser for LlamaIndex.

    Extracts text, tables, and speaker notes from PPTX files.
    Returns one Document per slide with rich metadata.
    """

    def __init__(
        self,
        raise_on_error: bool = False,
        describe_images: bool = False,
        image_dir: Optional[str] = None,
    ) -> None:
        """
        Initialize PptxReader.

        Args:
            raise_on_error: Whether to raise exceptions or continue on errors
            describe_images: Whether to use Claude Vision to describe slide images
            image_dir: Directory containing slide images (if describe_images=True)
        """
        self.raise_on_error = raise_on_error
        self.describe_images = describe_images
        self.image_dir = image_dir

    def load_data(
        self,
        file: Union[str, Path],
        extra_info: Optional[Dict] = None,
    ) -> List[Document]:
        """
        Parse PowerPoint file and extract content.

        Args:
            file: Path to the PowerPoint file
            extra_info: Additional metadata to include

        Returns:
            List of Documents (one per slide)
        """
        logger.info(f"Loading PPTX file: {file}")
        file_path = Path(file)

        try:
            # Open presentation
            prs = Presentation(str(file_path))

            # Process each slide
            docs = []
            for slide_num, slide in enumerate(prs.slides, start=1):
                try:
                    # Extract text from all shapes
                    slide_text_parts = []
                    tables_info = []

                    for shape in slide.shapes:
                        # Extract text from text shapes
                        if hasattr(shape, "text") and shape.text.strip():
                            slide_text_parts.append(shape.text.strip())

                        # Extract tables
                        if shape.has_table:
                            table = shape.table
                            table_data = []
                            for row in table.rows:
                                row_data = [cell.text.strip() for cell in row.cells]
                                table_data.append(" | ".join(row_data))
                            if table_data:
                                tables_info.append("\n".join(table_data))

                    # Combine all text
                    slide_text = "\n".join(slide_text_parts)

                    # Extract speaker notes
                    notes = ""
                    try:
                        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                            notes = slide.notes_slide.notes_text_frame.text.strip()
                    except Exception as e:
                        logger.warning(f"Could not extract notes from slide {slide_num}: {e}")

                    # Optionally add image description
                    image_description = ""
                    if self.describe_images and self.image_dir:
                        try:
                            # Look for slide image
                            from utils import describe_image
                            file_name_base = Path(file_path).stem.replace(' ', '_')
                            image_path = Path(self.image_dir) / f"{file_name_base}_{slide_num-1:04d}.png"

                            if image_path.exists():
                                with open(image_path, 'rb') as img_file:
                                    image_content = img_file.read()
                                image_description = describe_image(image_content)
                                logger.info(f"Slide {slide_num}: Added visual description")
                        except Exception as e:
                            logger.warning(f"Slide {slide_num}: Could not describe image: {e}")

                    # Build comprehensive content
                    content_parts = []

                    if slide_text:
                        content_parts.append(f"Slide {slide_num} content:\n{slide_text}")

                    if tables_info:
                        content_parts.append(f"\nTables:\n" + "\n\n".join(tables_info))

                    if image_description:
                        content_parts.append(f"\nVisual content:\n{image_description}")

                    if notes:
                        content_parts.append(f"\nSpeaker notes:\n{notes}")

                    content = "\n\n".join(content_parts) if content_parts else f"Slide {slide_num} (no text content)"

                    # Create metadata
                    metadata = {
                        "file_path": str(file_path),
                        "file_name": file_path.name,
                        "page_label": slide_num,
                        "slide_number": slide_num,
                        "has_tables": len(tables_info) > 0,
                        "has_notes": bool(notes),
                        "source": f"{file_path.name} - Slide {slide_num}",
                    }

                    if extra_info:
                        metadata.update(extra_info)

                    # Create document
                    doc = Document(
                        text=content,
                        metadata=metadata,
                    )
                    docs.append(doc)

                except Exception as e:
                    error_msg = f"Error processing slide {slide_num}: {e}"
                    logger.error(error_msg)
                    if self.raise_on_error:
                        raise
                    # Create error document
                    docs.append(Document(
                        text=f"Slide {slide_num}: Error during extraction",
                        metadata={
                            "file_path": str(file_path),
                            "slide_number": slide_num,
                            "error": str(e),
                        }
                    ))

            logger.info(f"Successfully extracted {len(docs)} slides from {file_path.name}")
            return docs

        except Exception as e:
            error_msg = f"Failed to open PowerPoint file {file}: {e}"
            logger.error(error_msg)
            if self.raise_on_error:
                raise
            return []
